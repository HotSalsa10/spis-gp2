"""Patch the SPIS presentation against verified source-of-truth values.

Reads the original PowerPoint, applies a list of paragraph-level text
replacements anchored to specific slide indices, and writes a new file.

Each replacement is logged so we can see which ones landed and which were
skipped (because the text was edited in the original).
"""

from __future__ import annotations

import shutil
import sys
from pathlib import Path

from pptx import Presentation


SRC = Path(r"C:\Users\smsh3\Downloads\SPIS_Final_Presentation2pptx.pptx")
DST = Path(r"C:\Users\smsh3\Downloads\SPIS_Final_Presentation_FIXED.pptx")


# (slide_index_1based, old_paragraph_text_exact_or_substring, new_paragraph_text)
# We use substring matching on the *full paragraph text* (joined across runs)
# so we are robust to multi-run formatting.
# Nth-occurrence edits (slide_index, old, new, nth) — nth is 1-based.
# Used when the same paragraph text appears multiple times on the same slide.
EDITS_NTH: list[tuple[int, str, str, int]] = [
    # Slide 24: first "6 tests" = Database/Catalog row -> 29
    (24, "6 tests", "29 tests", 1),
    # Slide 24 Test Types breakdown: Unit/Integration/Error Handling totals.
    # Each of these appears only once, so first-occurrence is fine — they're
    # under "Test Types" further down the slide.
    (24, "40 tests", "~100 tests", 1),
    (24, "27 tests", "~60 tests", 1),
    # "8 tests" appears TWICE on slide 24: first under XGBoost Forecaster
    # (which actually has 8 tests — leave alone), second under Error Handling.
    # Only update the SECOND occurrence.
    (24, "8 tests", "~22 tests", 2),
]


EDITS: list[tuple[int, str, str]] = [
    # ---- Slide 20: Tech stack — minor version bumps ----
    (20, "joblib 1.3",
         "joblib 1.5"),
    (20, "holidays 0.35",
         "holidays 0.50"),

    # ---- Slide 24 (second pass): module table ----
    (24, "Database Schema",
         "Database / Catalog"),  # rebadge to reflect added Phase 9 scope

    # ---- Slide 15: Database Design — Four Tables → Eight Tables ----
    (15, "Database Design — Four Tables",
         "Database Design — Eight Tables"),
    (15, "A SQLite schema with two reference tables (seeded once) and two operational tables (updated at runtime).",
         "A SQLite schema with the four primary tables shown below plus four Phase-9 workflow tables (inventory_batches, alerts, suppliers, purchase_orders)."),

    # ---- Slide 16: GridSearch combos 512 → 128 ----
    (16, "XGBoost + GridSearch (512 combos, TimeSeriesSplit)",
         "XGBoost + GridSearch (128 combos, TimeSeriesSplit)"),

    # ---- Slide 21: Training XGBoost ----
    (21, "Drop NaN lag rows (first 7 days per drug)",
         "Drop NaN lag rows (first ~365 days per drug, due to lag_365)"),
    (21, "GridSearchCV across 512 parameter combinations",
         "GridSearchCV across 128 parameter combinations"),
    (21, "n_estimators = 500",
         "n_estimators = 800"),
    (21, "learning_rate = 0.1",
         "learning_rate = 0.03"),
    (21, "min_child_weight = 5",
         "min_child_weight = 1"),

    # ---- Slide 24: Testing 75 → 182, 7 files → 14 files ----
    (24, "Testing — 75 Tests, 7 Modules, 80%+ Coverage",
         "Testing — 182 Tests, 14 Modules, 80%+ Coverage"),
    (24, "75",
         "182"),
    (24, "7",
         "14"),
    # The "12s full run" estimate updates because 182 tests run longer
    (24, "~12s",
         "~25s"),
    # Module-level counts that changed
    (24, "Database Schema",
         "Database Schema"),  # placeholder — actual count update below
    # The number "6" for Database Schema tests → 29
    # (handled via slide-paragraph search by neighbouring text)

    # ---- Slide 28: Objectives — outdated metadata ----
    (28, "SQLite, 4 tables, 57 drugs, 424k records",
         "SQLite, 8 tables, 57 drugs, 424k records"),
    (28, "4 sections, cached, missing-file safe-guard",
         "Multi-page (8 pages), cached, missing-file safe-guard"),
    (28, "75 tests, 7 files, 80%+ coverage",
         "182 tests, 14 files, 80%+ coverage"),

    # ---- Slide 29: Limitations — fix the "no batch-level expiry" claim ----
    (29, "No batch-level expiry",
         "Batch expiry seeded, not live"),
    (29, "Public dataset has no expiry dates, so SKU-level expiry alerts are not yet possible.",
         "Public sales dataset has no batch expiry dates. We seeded demo batches into inventory_batches and built the expiry advisor; a live POS feed would supply real batch metadata."),

    # ---- Slide 31: Conclusion footer "75 tests" → "182 tests" ----
    (31, "Built in one semester  •  75 tests  •  80%+ coverage  •  Open source",
         "Built across two semesters  •  182 tests  •  80%+ coverage  •  Open source"),

    # ---- Slide 37: Q&A #5 reframe (expiry tier exists in form of advisor) ----
    (37, "Why is there no expiry-risk tier in the current system?",
         "How does SPIS handle drug-batch expiry today?"),
    (37, "It's a dataset limitation, not a design flaw — and the architecture is ready for it.",
         "An expiry advisor is implemented and runs over a seeded inventory_batches table; production needs a live POS feed for real batch data."),
    (37, "The data gap",
         "What is implemented"),
    (37, "The Kaggle Pharma Sales dataset records transaction quantities but does not include batch-level expiry dates. Without per-batch shelf-life data, we cannot tell which units in current stock will expire next.",
         "spis/models/expiry_advisor.py implements a two-factor discount classifier on days_to_expiry and risk_ratio. inventory_batches stores expiry_date and unit_cost per batch. The Expiry Offers dashboard page surfaces the recommended action (none/discount/return/write_off) for each batch."),
    (37, "What's already in place",
         "What would extend it to production"),
    (37, "Requirements (Ch 3) already specify expiry-risk classification as a target tier",
         "Live POS feed to register batches at receive time (currently CSV-seeded)"),
    (37, "RiskAssessment dataclass can be extended with batch metadata without breaking existing code",
         "Per-drug (not per-ATC) sell-through forecasting to predict batch consumption"),
    (37, "Configurable threshold framework supports adding an expiry-window parameter",
         "Probabilistic time-to-expiry confidence intervals on each batch"),
    (37, "Dashboard's color-coded alert pattern naturally extends to a fifth tier",
         "Integration with the alert engine so write-off events page operators automatically"),
]


def paragraph_full_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs)


def replace_paragraph_text(paragraph, new_text: str) -> None:
    """Set paragraph text to *new_text*, preserving the first run's formatting.

    Strategy: keep the first run, set its text to new_text, blank out
    subsequent runs."""
    runs = paragraph.runs
    if not runs:
        return
    runs[0].text = new_text
    for r in runs[1:]:
        r.text = ""


def apply_edits(prs: Presentation, edits: list[tuple[int, str, str]]) -> tuple[int, int]:
    landed = 0
    skipped = 0
    # build lookup: slide_idx -> list of (old, new)
    by_slide: dict[int, list[tuple[str, str]]] = {}
    for s, o, n in edits:
        by_slide.setdefault(s, []).append((o, n))

    for slide_idx, edits_here in by_slide.items():
        if slide_idx < 1 or slide_idx > len(prs.slides):
            print(f"  [skip] slide {slide_idx} out of range")
            continue
        slide = prs.slides[slide_idx - 1]
        for old, new in edits_here:
            found = False
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    full = paragraph_full_text(paragraph)
                    if full == old:
                        replace_paragraph_text(paragraph, new)
                        landed += 1
                        found = True
                        print(f"  [ok ] slide {slide_idx}: '{old[:60]}' -> '{new[:60]}'")
                        break
                if found:
                    break
            # try tables too if not found
            if not found:
                for shape in slide.shapes:
                    if not shape.has_table:
                        continue
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                full = paragraph_full_text(paragraph)
                                if full == old:
                                    replace_paragraph_text(paragraph, new)
                                    landed += 1
                                    found = True
                                    print(f"  [ok ] slide {slide_idx} (table): '{old[:60]}' -> '{new[:60]}'")
                                    break
                            if found:
                                break
                        if found:
                            break
                    if found:
                        break
            if not found:
                skipped += 1
                print(f"  [MISS] slide {slide_idx}: could not find '{old[:80]}'")
    return landed, skipped


def apply_nth_edits(prs: Presentation, edits: list[tuple[int, str, str, int]]) -> tuple[int, int]:
    """Replace the Nth occurrence (1-based) of *old* on each slide."""
    landed = 0
    skipped = 0
    by_slide: dict[int, list[tuple[str, str, int]]] = {}
    for s, o, n, nth in edits:
        by_slide.setdefault(s, []).append((o, n, nth))
    for slide_idx, items in by_slide.items():
        if slide_idx < 1 or slide_idx > len(prs.slides):
            continue
        slide = prs.slides[slide_idx - 1]

        def find_and_replace_nth(old: str, new: str, nth: int) -> bool:
            seen = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph_full_text(paragraph) == old:
                            seen += 1
                            if seen == nth:
                                replace_paragraph_text(paragraph, new)
                                return True
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                if paragraph_full_text(paragraph) == old:
                                    seen += 1
                                    if seen == nth:
                                        replace_paragraph_text(paragraph, new)
                                        return True
            return False

        for old, new, nth in items:
            if find_and_replace_nth(old, new, nth):
                landed += 1
                print(f"  [ok ] slide {slide_idx} (nth={nth}): '{old}' -> '{new}'")
            else:
                skipped += 1
                print(f"  [MISS] slide {slide_idx} (nth={nth}): '{old}'")
    return landed, skipped


def main() -> int:
    if not SRC.exists():
        print(f"source not found: {SRC}")
        return 1
    shutil.copy(SRC, DST)
    prs = Presentation(str(DST))
    print(f"loaded {len(prs.slides)} slides from {DST.name}")
    landed, skipped = apply_edits(prs, EDITS)
    print()
    landed2, skipped2 = apply_nth_edits(prs, EDITS_NTH)
    prs.save(str(DST))
    total_landed = landed + landed2
    total_skipped = skipped + skipped2
    print(f"\nresult: {total_landed} edits applied, {total_skipped} skipped")
    print(f"saved to: {DST}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
